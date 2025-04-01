import os
from flask import Flask, request, jsonify
from pptx import Presentation

app = Flask(__name__)

@app.route('/')
def home():
    return "SmartArt API is running!"

@app.route('/smartart', methods=['POST'])
def apply_smartart():
    try:
        # Save uploaded file
        file = request.files['file']
        filepath = f"/tmp/{file.filename}"
        file.save(filepath)

        # Open the PPTX
        prs = Presentation(filepath)

        # Example formatting: change title text to SmartArt style
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        paragraph.level = 0  # Change indentation level if needed

        # Save the modified file
        new_path = f"/tmp/modified_{file.filename}"
        prs.save(new_path)

        return jsonify({"message": "SmartArt formatting applied", "output": new_path})

    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
